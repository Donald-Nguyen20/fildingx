# Funtion/rag_extract.py
import json
import csv

def extract_content(file_path: str) -> str:
    fp = (file_path or "").lower()

    # TEXT
    if fp.endswith(".txt") or fp.endswith(".md"):
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()

    # PDF (fitz / pymupdf)
    if fp.endswith(".pdf"):
        import fitz
        with fitz.open(file_path) as pdf:
            return "\n".join(page.get_text("text") for page in pdf)

    # DOCX
    if fp.endswith(".docx"):
        from docx import Document
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs)

    # XLSX
    if fp.endswith(".xlsx"):
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True, data_only=True)
        out = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                out.append("\t".join("" if v is None else str(v) for v in row))
        return "\n".join(out)

    # PPTX (rút text basic)
    if fp.endswith(".pptx"):
        from pptx import Presentation
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        return "\n".join(texts)

    # CSV
    if fp.endswith(".csv"):
        rows = []
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            for r in reader:
                rows.append(", ".join(r))
        return "\n".join(rows)

    # JSON
    if fp.endswith(".json"):
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            data = json.load(f)
        return json.dumps(data, ensure_ascii=False, indent=2)

    # XML (fallback: đọc raw)
    if fp.endswith(".xml"):
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()

    # HTML
    if fp.endswith(".html") or fp.endswith(".htm"):
        from bs4 import BeautifulSoup
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            soup = BeautifulSoup(f, "html.parser")
        return soup.get_text("\n")

    return ""
